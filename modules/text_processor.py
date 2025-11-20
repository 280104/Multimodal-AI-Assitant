import os
import tempfile
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import markdown
#import whisper
from langchain.text_splitter import RecursiveCharacterTextSplitter

# --- Whisper model (load once) ---
whisper_model = whisper.load_model("base")


# --- Text extraction functions ---
def extract_text_from_pdf(file_obj):
    # Save uploaded file to temporary location
    if hasattr(file_obj, 'read'):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            tmp.write(file_obj.read())
            tmp_path = tmp.name
    else:
        tmp_path = file_obj
    
    text = ""
    with fitz.open(tmp_path) as doc:
        for page in doc:
            text += page.get_text()
    
    # Clean up temp file
    if hasattr(file_obj, 'read'):
        os.remove(tmp_path)
    
    return text


def extract_text_from_docx(file_obj):
    if hasattr(file_obj, 'read'):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_obj.read())
            tmp_path = tmp.name
        doc = docx.Document(tmp_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        os.remove(tmp_path)
        return text
    else:
        doc = docx.Document(file_obj)
        return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_pptx(file_obj):
    if hasattr(file_obj, 'read'):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(file_obj.read())
            tmp_path = tmp.name
        prs = Presentation(tmp_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        os.remove(tmp_path)
        return text
    else:
        prs = Presentation(file_obj)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text


def extract_text_from_md(file_obj):
    if hasattr(file_obj, 'read'):
        md_content = file_obj.read().decode('utf-8')
    else:
        with open(file_obj, 'r', encoding='utf-8') as f:
            md_content = f.read()
    return markdown.markdown(md_content)


def extract_text_from_txt(file_obj):
    if hasattr(file_obj, 'read'):
        return file_obj.read().decode('utf-8')
    else:
        with open(file_obj, 'r', encoding='utf-8') as f:
            return f.read()


def extract_text_from_audio(audio_obj):
    if hasattr(audio_obj, 'read'):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp:
            tmp.write(audio_obj.read())
            tmp_path = tmp.name
        result = whisper_model.transcribe(tmp_path)
        os.remove(tmp_path)
        return result["text"]
    else:
        result = whisper_model.transcribe(audio_obj)
        return result["text"]

# --- Master extractor ---
def extract_text(file_or_link):
    # Uploaded file only (YouTube removed)
    if hasattr(file_or_link, "name"):
        ext = os.path.splitext(file_or_link.name)[1].lower()
        if ext == ".pdf":
            return extract_text_from_pdf(file_or_link)
        elif ext == ".docx":
            return extract_text_from_docx(file_or_link)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_or_link)
        elif ext == ".md":
            return extract_text_from_md(file_or_link)
        elif ext == ".txt":
            return extract_text_from_txt(file_or_link)
        elif ext in [".mp3", ".mp4"]:
            return extract_text_from_audio(file_or_link)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    else:
        raise ValueError("Invalid input. Please upload a file.")
# --- Master extractor ---
def extract_text(file_or_link):
    # Uploaded file
    if hasattr(file_or_link, "name"):
        ext = os.path.splitext(file_or_link.name)[1].lower()
        if ext == ".pdf":
            return extract_text_from_pdf(file_or_link)
        elif ext == ".docx":
            return extract_text_from_docx(file_or_link)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_or_link)
        elif ext == ".md":
            return extract_text_from_md(file_or_link)
        elif ext == ".txt":
            return extract_text_from_txt(file_or_link)
        elif ext in [".mp3", ".mp4"]:
            return extract_text_from_audio(file_or_link)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    # YouTube link
   

# --- Text splitter ---
def split_text(text, chunk_size=500, chunk_overlap=50):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    return splitter.split_text(text)
